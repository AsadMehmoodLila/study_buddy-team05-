import os
import json
import io
import re
from openai import OpenAI
from dotenv import load_dotenv

try:
    from PyPDF2 import PdfReader
except ImportError:  # pragma: no cover - handled at runtime in the UI
    PdfReader = None

try:
    from pptx import Presentation
except ImportError:  # pragma: no cover - handled at runtime in the UI
    Presentation = None

# Load environment variables from .env file
load_dotenv()

MAX_CONTEXT_CHARS = 18000


def _missing_key_message():
    return "OPENROUTER_API_KEY is missing. Add it to your .env file and restart Streamlit."


def _get_providers():
    return [
        {
            "name": "OpenRouter",
            "api_key": os.getenv("OPENROUTER_API_KEY"),
            "base_url": "https://openrouter.ai/api/v1",
            "model": os.getenv("OPENROUTER_MODEL", "openrouter/free"),
        },
        {
            "name": "Groq",
            "api_key": os.getenv("GROQ_API_KEY"),
            "base_url": "https://api.groq.com/openai/v1",
            "model": "llama-3.3-70b-versatile",
        },
        {
            "name": "Google Gemini",
            "api_key": os.getenv("GEMINI_API_KEY"),
            "base_url": "https://generativelanguage.googleapis.com/v1beta/openai/",
            "model": "gemini-1.5-flash",
        }
    ]


def _extract_topic(user_prompt):
    match = re.search(r"'(.*?)'", user_prompt)
    if match:
        return match.group(1)
    match_on = re.search(r"on: '(.*?)'", user_prompt)
    if match_on:
        return match_on.group(1)
    return "the uploaded documents"


def _mock_notes(topic):
    return {
        "notes": (
            f"Mock Offline Mode (Apology: All API keys failed or limits reached):\n\n"
            f"Here are simulated study notes for '{topic}':\n\n"
            f"1. Core Concepts: The study of {topic} focuses on identifying primary elements, "
            "analyzing key variables, and understanding logical frameworks.\n\n"
            f"2. Methodologies: Implementing techniques related to {topic} involves structured, "
            "repeatable procedures to ensure reproducible and reliable results.\n\n"
            "3. Best Practices: Students should utilize active recall, spaced repetition, "
            "and continuous testing to master these concepts."
        ),
        "points": [
            f"Understand the main definitions and context of {topic}.",
            f"Identify primary techniques, processes, or models.",
            f"Apply structured analysis to evaluate practical use-cases."
        ],
        "summary": f"In summary, mastering {topic} requires a structured combination of theoretical study and practical exercises."
    }


def _mock_flashcards(topic):
    return {
        "flashcards": [
            {"question": f"What is the primary definition of {topic}?", "answer": f"It refers to the core concepts and frameworks defined under {topic}."},
            {"question": f"Why is {topic} important in study planning?", "answer": "It serves as a fundamental building block for mastering advanced materials."},
            {"question": f"What is a common challenge when studying {topic}?", "answer": "Synthesizing large volumes of text into structured notes without active recall."},
            {"question": f"Name a key technique to study {topic} efficiently.", "answer": "Breaking down the materials into smaller components and testing yourself with flashcards."},
            {"question": f"How can you evaluate your progress in {topic}?", "answer": "By taking regular quizzes and analyzing your performance over time."}
        ]
    }


def _mock_quiz(topic, user_prompt):
    return {
        "quiz": [
            {
                "question": f"Which of the following is a primary objective of studying {topic}?",
                "options": [
                    "To optimize information retention and understanding",
                    "To completely avoid structured review sessions",
                    "To memorize facts without understanding their context",
                    "To rely solely on passive reading before exams"
                ],
                "answer": "To optimize information retention and understanding"
            },
            {
                "question": f"What is a key benefit of using active recall for {topic}?",
                "options": [
                    "It strengthens memory retrieval paths compared to passive reading",
                    "It guarantees 100% exam scores without study effort",
                    "It replaces the need for understanding the basics",
                    "It is only useful for advanced topics"
                ],
                "answer": "It strengthens memory retrieval paths compared to passive reading"
            },
            {
                "question": f"How should you structure study materials for {topic}?",
                "options": [
                    "Rank files logically by starting with foundational topics first",
                    "Study advanced topics before learning basic terminology",
                    "Complete tasks in random sequence order",
                    "Avoid using a timeline or study planner"
                ],
                "answer": "Rank files logically by starting with foundational topics first"
            }
        ]
    }


def _mock_roadmap(user_prompt):
    files = re.findall(r"FILE:\s*(.*?)\n", user_prompt)
    if not files:
        files = ["Document_1.pdf", "Document_2.pdf"]
    roadmap = []
    for idx, filename in enumerate(files, start=1):
        roadmap.append({
            "step": idx,
            "file": filename,
            "reason": f"Establish a solid understanding of foundations by reviewing {filename}.",
            "focus": f"Take comprehensive notes and review key concepts in {filename}."
        })
    return {"roadmap": roadmap}


def _mock_chat(user_prompt):
    context_match = re.search(r"Uploaded context:\s*(.*?)\s*Student question:", user_prompt, re.DOTALL)
    question_match = re.search(r"Student question:\s*(.*)", user_prompt, re.DOTALL)
    
    context = context_match.group(1).strip() if context_match else ""
    question = question_match.group(1).strip() if question_match else ""
    
    if context and question:
        words = [w.lower() for w in re.findall(r"\b\w{4,}\b", question)]
        sentences = re.split(r"(?<=[.!?])\s+", context)
        
        matches = []
        for sentence in sentences:
            score = sum(1 for word in words if word in sentence.lower())
            if score > 0:
                matches.append((score, sentence))
        
        if matches:
            matches.sort(key=lambda x: x[0], reverse=True)
            best_sentences = [m[1] for m in matches[:3]]
            joined_sentences = " ".join(best_sentences)
            return (
                f"Mock Offline Mode (Apology: All API keys failed or limits reached):\n\n"
                f"[Simulated response based on document text]:\n"
                f"{joined_sentences}\n\n"
                f"(Source Grounding: Extracted from matching sentences in context)"
            )
            
    return (
        f"Mock Offline Mode (Apology: All API keys failed or limits reached):\n\n"
        f"I could not find a specific answer for '{question}' in the uploaded documents. "
        "Please check if all API services are online or try rephrasing your question."
    )


def _mock_summary(user_prompt):
    doc_name_match = re.search(r"Document name:\s*(.*?)\n", user_prompt)
    doc_name = doc_name_match.group(1).strip() if doc_name_match else "Document"
    
    context_match = re.search(r"Source document:\s*(.*)", user_prompt, re.DOTALL)
    context = context_match.group(1).strip() if context_match else ""
    
    sentences = re.split(r"(?<=[.!?])\s+", context)
    valid_sentences = [s for s in sentences if len(s.strip()) > 10]
    
    if len(valid_sentences) >= 2:
        extracted = " ".join(valid_sentences[:3])
        return f"Mock Offline Mode Summary for '{doc_name}': {extracted}"
    else:
        return f"Mock Offline Mode Summary: This document covers key concepts and structures relating to '{doc_name}'."


def _local_mock_fallback(system_prompt, user_prompt, json_mode, error_context=None):
    print(f"API Fallback system activated. Error occurred: {error_context}. Running local mock fallback.")
    if json_mode:
        if "flashcards" in user_prompt.lower():
            topic = _extract_topic(user_prompt)
            return _mock_flashcards(topic)
        elif "quiz" in user_prompt.lower():
            topic = _extract_topic(user_prompt)
            return _mock_quiz(topic, user_prompt)
        elif "roadmap" in user_prompt.lower():
            return _mock_roadmap(user_prompt)
        else:
            topic = _extract_topic(user_prompt)
            return _mock_notes(topic)
    else:
        if "student question:" in user_prompt.lower():
            return _mock_chat(user_prompt)
        elif "document name:" in user_prompt.lower():
            return _mock_summary(user_prompt)
        else:
            return (
                "Mock Offline Mode (Apology: All API keys failed or limits reached):\n\n"
                "An error occurred with all configured API providers. Here is a simulated response to keep the application running."
            )


def _call_llm_with_fallback(system_prompt, user_prompt, json_mode=False):
    providers = _get_providers()
    last_error = None
    
    for provider in providers:
        api_key = provider["api_key"]
        if not api_key or api_key == "missing-key" or api_key.startswith("your_actual"):
            continue
        
        try:
            prov_client = OpenAI(
                base_url=provider["base_url"],
                api_key=api_key
            )
            
            kwargs = {
                "model": provider["model"],
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_prompt}
                ]
            }
            if json_mode:
                kwargs["response_format"] = {"type": "json_object"}
            
            response = prov_client.chat.completions.create(**kwargs)
            content = response.choices[0].message.content
            if not content:
                raise ValueError("Received empty content from provider API.")
            
            if json_mode:
                return json.loads(content)
            else:
                return content.strip()
                
        except Exception as e:
            last_error = f"[{provider['name']} failed: {str(e)}]"
            
    return _local_mock_fallback(system_prompt, user_prompt, json_mode, error_context=last_error)


def _json_completion(system_prompt, user_prompt, fallback):
    try:
        return _call_llm_with_fallback(system_prompt, user_prompt, json_mode=True)
    except Exception as e:
        return fallback(str(e))


def _text_completion(system_prompt, user_prompt):
    try:
        return _call_llm_with_fallback(system_prompt, user_prompt, json_mode=False)
    except Exception as e:
        return f"AI request failed: {str(e)}"


def _clean_text(text):
    """Normalizes whitespace while keeping paragraph breaks readable."""
    text = text.replace("\x00", " ")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _truncate_context(context, max_chars=MAX_CONTEXT_CHARS):
    if not context:
        return ""
    if len(context) <= max_chars:
        return context
    return context[:max_chars] + "\n\n[Context truncated to fit the model window.]"


def extract_text_from_file(file_name, file_bytes):
    """Extracts text from PDF, text/markdown, and PowerPoint files."""
    extension = os.path.splitext(file_name.lower())[1]

    if extension == ".pdf":
        if PdfReader is None:
            return "", "Install PyPDF2 to extract PDF text."
        try:
            reader = PdfReader(io.BytesIO(file_bytes))
            page_text = []
            for page in reader.pages:
                page_text.append(page.extract_text() or "")
            return _clean_text("\n\n".join(page_text)), None
        except Exception as e:
            return "", f"Could not read PDF: {str(e)}"

    if extension in {".txt", ".md"}:
        for encoding in ("utf-8", "utf-16", "latin-1"):
            try:
                return _clean_text(file_bytes.decode(encoding)), None
            except UnicodeDecodeError:
                continue
        return "", "Could not decode text file."

    if extension in {".pptx"}:
        if Presentation is None:
            return "", "Install python-pptx to extract slide text."
        try:
            deck = Presentation(io.BytesIO(file_bytes))
            chunks = []
            for slide_number, slide in enumerate(deck.slides, start=1):
                chunks.append(f"Slide {slide_number}")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        chunks.append(shape.text)
            return _clean_text("\n".join(chunks)), None
        except Exception as e:
            return "", f"Could not read slides: {str(e)}"

    return "", "Unsupported file type."


def extract_documents(file_payloads):
    """Builds normalized document records from cached upload payloads."""
    documents = []
    for payload in file_payloads:
        file_name = payload["name"]
        text, error = extract_text_from_file(file_name, payload["bytes"])
        documents.append({
            "name": file_name,
            "type": payload.get("type", "unknown"),
            "text": text,
            "error": error,
            "char_count": len(text),
            "word_count": len(text.split())
        })
    return documents


def build_context_block(documents, max_chars=MAX_CONTEXT_CHARS):
    """Combines uploaded documents into a source-of-truth context block."""
    sections = []
    for doc in documents:
        if doc.get("text"):
            sections.append(f"DOCUMENT: {doc['name']}\n{doc['text']}")
    return _truncate_context("\n\n---\n\n".join(sections), max_chars=max_chars)

def generate_short_summary(document_name, document_text):
    """Creates a brief 2-3 line summary for individual file view."""
    context = _truncate_context(document_text, max_chars=14000)
    system_prompt = (
        "You are a careful study assistant. Summarize only the provided document. "
        "Do not add outside facts. Keep your summary extremely concise, strictly 2 to 3 lines maximum."
    )
    user_prompt = f"""
    Document name: {document_name}

    Create a brief 2-3 line summary capturing the core idea of this document.

    Source document:
    {context}
    """
    return _text_completion(system_prompt, user_prompt)

def generate_detailed_summary(document_name, document_text):
    """Creates a comprehensive 4-6 sentence summary for the combined view."""
    context = _truncate_context(document_text, max_chars=14000)
    system_prompt = (
        "You are a careful study assistant. Summarize the provided document comprehensively. "
        "Do not add outside facts. Provide a well-explained summary."
    )
    user_prompt = f"""
    Document name: {document_name}

    Create a clear and meaningful summary for this document. Explain the core concepts in about 4 to 6 sentences so the student can properly understand the main ideas and context.

    Source document:
    {context}
    """
    return _text_completion(system_prompt, user_prompt)


def generate_notes(topic, difficulty="Intermediate", source_context=None):
    """Generates structured study notes for a given topic."""
    context_instruction = ""
    if source_context:
        context_instruction = f"""
        Use the uploaded document context as the only source of truth.
        If the answer is not supported by the context, say that the uploaded documents do not contain enough information.

        Uploaded context:
        {_truncate_context(source_context)}
        """

    prompt = f"""
    Create comprehensive study notes for: '{topic or "the uploaded documents"}'.
    Target difficulty level: {difficulty}.
    {context_instruction}
    
    Output strictly in the following JSON format:
    {{
        "notes": "A detailed explanation of the topic.",
        "points": ["Key point 1", "Key point 2", "Key point 3"],
        "summary": "A brief one-paragraph summary."
    }}
    """
    def fallback(error):
        return {
            "notes": f"Error generating notes: {error}",
            "points": ["Error retrieving key points."],
            "summary": "Please try again later."
        }

    return _json_completion(
        "You are a helpful AI study assistant that strictly outputs valid JSON.",
        prompt,
        fallback
    )

def generate_flashcards(topic, source_context=None):
    """Generates a list of flashcards for active recall."""
    context_instruction = ""
    if source_context:
        context_instruction = f"""
        Use only this uploaded document context:
        {_truncate_context(source_context)}
        """

    prompt = f"""
    Create 5 educational flashcards for: '{topic or "the uploaded documents"}'.
    {context_instruction}
    
    Output strictly in the following JSON format:
    {{
        "flashcards": [
            {{"question": "Question 1", "answer": "Answer 1"}},
            {{"question": "Question 2", "answer": "Answer 2"}}
        ]
    }}
    """

    def fallback(error):
        return {"flashcards": [{"question": "API Error", "answer": error}]}

    data = _json_completion(
        "You are a helpful AI study assistant that strictly outputs valid JSON.",
        prompt,
        fallback
    )
    if isinstance(data, dict):
        return data.get("flashcards", [])
    return []

def generate_quiz(topic, source_context=None):
    """Generates a multiple-choice quiz."""
    context_instruction = ""
    if source_context:
        context_instruction = f"""
        Use the uploaded document context as the only source of truth.
        Avoid questions that require outside knowledge.

        Uploaded context:
        {_truncate_context(source_context)}
        """

    prompt = f"""
    Create a 10-question multiple-choice quiz on: '{topic or "the uploaded documents"}'.
    Each question must have exactly 4 options, and one correct answer that perfectly matches one of the options.
    {context_instruction}
    
    Output strictly in the following JSON format:
    {{
        "quiz": [
            {{
                "question": "What is...?",
                "options": ["Option A", "Option B", "Option C", "Option D"],
                "answer": "Option B"
            }}
        ]
    }}
    """

    def fallback(error):
        return {"quiz": [{
            "question": "Quiz generation failed.",
            "options": ["Check API key", "Check model", "Try again", "Use shorter context"],
            "answer": error
        }]}

    data = _json_completion(
        "You are a helpful AI study assistant that strictly outputs valid JSON.",
        prompt,
        fallback
    )
    if isinstance(data, dict):
        return data.get("quiz", [])
    return []

def chat_with_documents(question, source_context):
    """Answers a student question strictly from uploaded document context."""
    system_prompt = (
        "You are Study Buddy's document chat tutor. Answer strictly from the uploaded context. "
        "If the context does not support the answer, say: 'I could not find that in the uploaded documents.' "
        "Keep answers concise and cite the document section or filename when possible."
    )
    user_prompt = f"""
    Uploaded context:
    {_truncate_context(source_context)}

    Student question:
    {question}
    """
    return _text_completion(system_prompt, user_prompt)


def generate_study_roadmap(summaries):
    """Ranks uploaded files into a logical study sequence."""
    joined = "\n\n".join(
        f"FILE: {item['name']}\nSUMMARY:\n{item['summary']}"
        for item in summaries
        if item.get("summary")
    )
    prompt = f"""
    You are sequencing course materials for a student.
    Use these document summaries to choose the most logical reading order.
    Put fundamentals before advanced or applied material.

    Return strictly this JSON format:
    {{
        "roadmap": [
            {{
                "step": 1,
                "file": "filename.pdf",
                "reason": "Why this file comes here.",
                "focus": "What the student should pay attention to."
            }}
        ]
    }}

    Summaries:
    {_truncate_context(joined, max_chars=16000)}
    """

    def fallback(error):
        return {"roadmap": [{"step": 1, "file": "Roadmap unavailable", "reason": error, "focus": "Try again after checking your API key and model."}]}

    data = _json_completion(
        "You are a careful curriculum planner that strictly outputs valid JSON.",
        prompt,
        fallback
    )
    if isinstance(data, dict):
        return data.get("roadmap", [])
    return []

def save_progress(topic, score):
    """Saves user quiz scores locally to a JSON file."""
    try:
        with open("progress.json", "r") as f:
            data = json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        data = {"topics": [], "scores": []}

    data["topics"].append(topic)
    data["scores"].append(score)

    with open("progress.json", "w") as f:
        json.dump(data, f, indent=4)

def load_progress():
    """Loads user progress data."""
    try:
        with open("progress.json", "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {"topics": [], "scores": []}
    
# --- PRODUCTIVITY ENGINE LOGIC ---

def load_productivity():
    """Loads tasks and pomodoro sessions from local storage."""
    try:
        with open("productivity.json", "r") as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        # Default structure if file doesn't exist
        return {"tasks": [], "pomodoro_sessions": 0}

def save_productivity(data):
    """Saves productivity data to local storage."""
    with open("productivity.json", "w") as f:
        json.dump(data, f, indent=4)

def normalize_task_priorities():
    """Ensures active tasks have sequential priorities starting at 1."""
    data = load_productivity()
    active_tasks = [task for task in data["tasks"] if not task.get("completed")]
    active_tasks.sort(key=lambda task: (task.get("priority", task.get("id", 0)), task.get("id", 0)))

    changed = False
    for index, task in enumerate(active_tasks, start=1):
        if task.get("priority") != index:
            task["priority"] = index
            changed = True

    for task in data["tasks"]:
        if task.get("completed") and task.get("priority") is not None:
            task["priority"] = None
            changed = True

    if changed:
        save_productivity(data)
    return data

def add_task(name, is_important=False, is_urgent=False):
    """Adds a new task to the productivity tracker."""
    data = normalize_task_priorities()
    active_count = len([task for task in data["tasks"] if not task.get("completed")])
    new_task = {
        "id": len(data["tasks"]) + 1,
        "name": name,
        "important": is_important,
        "urgent": is_urgent,
        "priority": active_count + 1,
        "completed": False
    }
    data["tasks"].append(new_task)
    save_productivity(data)

def update_task_priority(task_id, new_priority):
    """Moves an active task to a new priority and shifts the others."""
    data = normalize_task_priorities()
    active_tasks = [task for task in data["tasks"] if not task.get("completed")]
    active_tasks.sort(key=lambda task: (task.get("priority", task.get("id", 0)), task.get("id", 0)))

    target = next((task for task in active_tasks if task["id"] == task_id), None)
    if target is None:
        return

    active_tasks = [task for task in active_tasks if task["id"] != task_id]
    insert_at = max(0, min(int(new_priority) - 1, len(active_tasks)))
    active_tasks.insert(insert_at, target)

    for index, task in enumerate(active_tasks, start=1):
        task["priority"] = index

    save_productivity(data)

def complete_task(task_id):
    """Marks a specific task as completed."""
    data = load_productivity()
    for task in data["tasks"]:
        if task["id"] == task_id:
            task["completed"] = True
            task["priority"] = None
            break
    save_productivity(data)
    normalize_task_priorities()

def log_pomodoro_session():
    """Increments the total number of completed Pomodoro sessions."""
    data = load_productivity()
    data["pomodoro_sessions"] += 1
    save_productivity(data)
